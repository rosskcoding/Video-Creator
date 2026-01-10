"""
PPTX to PNG Converter
Uses LibreOffice for rendering and python-pptx for speaker notes extraction
"""
import asyncio
import hashlib
import os
import shutil
import tempfile
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches


class PPTXConverter:
    """Convert PPTX to PNG images and extract speaker notes"""
    
    def __init__(self, dpi: int = 150):
        """
        Args:
            dpi: Resolution for PNG export. 150 DPI gives ~1920px for 16:9 slides
        """
        self.dpi = dpi
    
    async def convert_to_png(
        self,
        pptx_path: Path,
        output_dir: Path,
    ) -> List[Path]:
        """
        Convert PPTX to PNG images using LibreOffice.
        
        Pipeline: PPTX → PDF → PNG (per page)
        
        Args:
            pptx_path: Path to input PPTX file
            output_dir: Directory for output PNG files
            
        Returns:
            List of paths to PNG files (in slide order)
        """
        pptx_path = Path(pptx_path)
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Create temp directory for intermediate files
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            # Step 1: Convert PPTX to PDF using LibreOffice
            pdf_path = temp_path / "presentation.pdf"
            await self._convert_to_pdf(pptx_path, temp_path)
            
            # Find the generated PDF (LibreOffice names it based on input)
            pdf_files = list(temp_path.glob("*.pdf"))
            if not pdf_files:
                raise RuntimeError("LibreOffice failed to generate PDF")
            pdf_path = pdf_files[0]
            
            # Step 2: Convert PDF to PNG pages
            png_paths = await self._pdf_to_png(pdf_path, output_dir)
        
        return png_paths
    
    async def _convert_to_pdf(self, pptx_path: Path, output_dir: Path) -> None:
        """Convert PPTX to PDF using LibreOffice headless"""
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_dir),
            str(pptx_path)
        ]
        
        process = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await process.communicate()
        
        if process.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {stderr.decode()}")
    
    async def _pdf_to_png(self, pdf_path: Path, output_dir: Path) -> List[Path]:
        """Convert PDF pages to PNG images using pdftoppm or ImageMagick"""
        # Try pdftoppm first (from poppler-utils), fallback to ImageMagick
        try:
            return await self._pdf_to_png_pdftoppm(pdf_path, output_dir)
        except FileNotFoundError:
            return await self._pdf_to_png_imagemagick(pdf_path, output_dir)
    
    async def _pdf_to_png_pdftoppm(self, pdf_path: Path, output_dir: Path) -> List[Path]:
        """Convert using pdftoppm (faster, better quality)"""
        output_prefix = output_dir / "slide"
        
        cmd = [
            "pdftoppm",
            "-png",
            "-r", str(self.dpi),
            str(pdf_path),
            str(output_prefix)
        ]
        
        process = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await process.communicate()
        
        if process.returncode != 0:
            raise RuntimeError(f"pdftoppm conversion failed: {stderr.decode()}")
        
        # pdftoppm creates files like: slide-1.png, slide-2.png, ...
        # Rename to: 001.png, 002.png, ...
        png_files = sorted(output_dir.glob("slide-*.png"))
        result = []
        
        for i, png_file in enumerate(png_files, 1):
            new_name = output_dir / f"{i:03d}.png"
            png_file.rename(new_name)
            result.append(new_name)
        
        return result
    
    async def _pdf_to_png_imagemagick(self, pdf_path: Path, output_dir: Path) -> List[Path]:
        """Convert using ImageMagick (fallback)"""
        cmd = [
            "convert",
            "-density", str(self.dpi),
            str(pdf_path),
            "-quality", "95",
            str(output_dir / "%03d.png")
        ]
        
        process = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await process.communicate()
        
        if process.returncode != 0:
            raise RuntimeError(f"ImageMagick conversion failed: {stderr.decode()}")
        
        # ImageMagick creates 000.png, 001.png, ...
        # Rename to 001.png, 002.png, ...
        png_files = sorted(output_dir.glob("*.png"))
        result = []
        
        for i, png_file in enumerate(png_files, 1):
            new_name = output_dir / f"{i:03d}.png"
            if png_file != new_name:
                png_file.rename(new_name)
            result.append(new_name)
        
        return result
    
    def extract_speaker_notes(self, pptx_path: Path) -> List[Optional[str]]:
        """
        Extract speaker notes from each slide.
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            List of notes text (or None if no notes) for each slide
        """
        pptx_path = Path(pptx_path)
        prs = Presentation(str(pptx_path))
        
        notes = []
        for slide in prs.slides:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text.strip()
                notes.append(notes_text if notes_text else None)
            else:
                notes.append(None)
        
        return notes
    
    def compute_slides_hash(self, pptx_path: Path) -> str:
        """
        Compute hash of PPTX structure for change detection.
        
        Uses slide count + slide dimensions + content hash
        """
        pptx_path = Path(pptx_path)
        prs = Presentation(str(pptx_path))
        
        # Collect structural info
        info_parts = [
            f"slides:{len(prs.slides)}",
            f"width:{prs.slide_width}",
            f"height:{prs.slide_height}",
        ]
        
        # Add hash of each slide's shapes (simplified)
        for i, slide in enumerate(prs.slides):
            shape_count = len(slide.shapes)
            info_parts.append(f"slide_{i}:shapes_{shape_count}")
        
        content = "|".join(info_parts)
        return hashlib.sha256(content.encode()).hexdigest()
    
    def compute_slide_hash(self, image_path: Path) -> str:
        """Compute hash of slide image for change detection"""
        with open(image_path, "rb") as f:
            return hashlib.sha256(f.read()).hexdigest()


# Singleton instance
pptx_converter = PPTXConverter()

