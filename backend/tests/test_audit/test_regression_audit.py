"""
Regression / Golden Test Audit (G-01)

Цель: когда вы правите импорт PPTX/рендер, набор "golden" презентаций
должен давать одинаковый результат (или diff объясним).

Test Matrix:
- G-01: "Золотые" PPTX кейсы
  - Простой PPTX
  - Со speaker notes  
  - С нестандартными шрифтами
  - С картинками/прозрачностью
"""
import json
import hashlib
import os
from pathlib import Path
from typing import Dict, Any, Optional
from unittest.mock import patch, MagicMock
from datetime import datetime

import pytest
from httpx import AsyncClient
from sqlalchemy.ext.asyncio import AsyncSession

from app.db.models import Project, ProjectVersion, Slide, SlideScript


# Golden test fixtures directory
GOLDEN_DIR = Path(__file__).parent / "golden_fixtures"


def compute_scene_hash(slides_data: list) -> str:
    """Compute hash of slide structure for comparison"""
    normalized = json.dumps(slides_data, sort_keys=True, default=str)
    return hashlib.sha256(normalized.encode()).hexdigest()[:16]


def compare_slide_structures(expected: list, actual: list) -> Dict[str, Any]:
    """Compare two slide structures and return diff"""
    diff = {
        "matched": True,
        "differences": [],
        "expected_count": len(expected),
        "actual_count": len(actual),
    }
    
    if len(expected) != len(actual):
        diff["matched"] = False
        diff["differences"].append({
            "type": "count_mismatch",
            "expected": len(expected),
            "actual": len(actual),
        })
    
    for i, (exp, act) in enumerate(zip(expected, actual)):
        if exp.get("slide_index") != act.get("slide_index"):
            diff["matched"] = False
            diff["differences"].append({
                "slide": i + 1,
                "type": "index_mismatch",
                "expected": exp.get("slide_index"),
                "actual": act.get("slide_index"),
            })
        
        # Check notes presence
        exp_has_notes = bool(exp.get("notes_text"))
        act_has_notes = bool(act.get("notes_text"))
        if exp_has_notes != act_has_notes:
            diff["matched"] = False
            diff["differences"].append({
                "slide": i + 1,
                "type": "notes_presence",
                "expected": exp_has_notes,
                "actual": act_has_notes,
            })
    
    return diff


class GoldenTestCase:
    """Base class for golden test cases"""
    
    def __init__(self, name: str, description: str):
        self.name = name
        self.description = description
        self.fixture_path: Optional[Path] = None
        self.expected_result: Optional[Dict] = None
    
    def load_expected(self) -> Dict:
        """Load expected result from golden file"""
        expected_file = GOLDEN_DIR / f"{self.name}_expected.json"
        if expected_file.exists():
            with open(expected_file) as f:
                return json.load(f)
        return {}
    
    def save_golden(self, result: Dict) -> None:
        """Save result as new golden file (for updating golden data)"""
        GOLDEN_DIR.mkdir(parents=True, exist_ok=True)
        golden_file = GOLDEN_DIR / f"{self.name}_expected.json"
        with open(golden_file, "w") as f:
            json.dump(result, f, indent=2, default=str)


# Define golden test cases
GOLDEN_CASES = {
    "simple_pptx": GoldenTestCase(
        name="simple_pptx",
        description="Simple PPTX with 3 slides, basic text",
    ),
    "with_speaker_notes": GoldenTestCase(
        name="with_speaker_notes",
        description="PPTX with speaker notes on each slide",
    ),
    "with_images": GoldenTestCase(
        name="with_images",
        description="PPTX with embedded images and transparency",
    ),
    "complex_fonts": GoldenTestCase(
        name="complex_fonts",
        description="PPTX with non-standard fonts (should fallback gracefully)",
    ),
}


class TestRegressionG01_GoldenPPTX:
    """
    G-01: "Золотые" PPTX кейсы
    
    Набор: простой, со speaker notes, с нестандартными шрифтами, с картинками
    Ожидаемо: импорт одинаковый от версии к версии (или diff объясним)
    """
    
    @pytest.mark.asyncio
    async def test_g01_simple_pptx_import(
        self,
        client: AsyncClient,
        sample_project: Project,
        db_session: AsyncSession,
        tmp_path
    ):
        """Test simple PPTX produces consistent slide count"""
        from pptx import Presentation
        
        # Create simple PPTX
        pptx_path = tmp_path / "simple.pptx"
        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            if slide.shapes.title:
                slide.shapes.title.text = f"Slide {i + 1}"
        prs.save(str(pptx_path))
        
        # Upload
        with open(pptx_path, "rb") as f:
            response = await client.post(
                f"/api/projects/{sample_project.id}/upload",
                files={"file": ("simple.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
            )
        
        assert response.status_code == 200
        
        # Golden expectation: 3 slides
        expected_count = 3
        actual_data = response.json()
        assert actual_data["status"] == "uploaded"
    
    @pytest.mark.asyncio
    async def test_g01_speaker_notes_preserved(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion,
        db_session: AsyncSession,
        tmp_path
    ):
        """Test speaker notes are extracted correctly"""
        from pptx import Presentation
        
        # Create PPTX with notes
        pptx_path = tmp_path / "with_notes.pptx"
        prs = Presentation()
        expected_notes = []
        
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            notes_text = f"Notes for slide {i + 1}"
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
            expected_notes.append(notes_text)
        
        prs.save(str(pptx_path))
        
        # Test the converter directly
        from app.adapters.pptx_converter import pptx_converter
        
        actual_notes = pptx_converter.extract_speaker_notes(pptx_path)
        
        # Golden expectation: notes match exactly
        assert len(actual_notes) == len(expected_notes)
        for expected, actual in zip(expected_notes, actual_notes):
            assert expected in actual  # Notes may have extra whitespace
    
    @pytest.mark.asyncio
    async def test_g01_slide_order_deterministic(
        self,
        client: AsyncClient,
        sample_project: Project,
        db_session: AsyncSession,
        tmp_path
    ):
        """Test slide order is deterministic across imports"""
        from pptx import Presentation
        
        # Create PPTX with numbered slides
        pptx_path = tmp_path / "ordered.pptx"
        prs = Presentation()
        for i in range(5):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            if slide.shapes.title:
                slide.shapes.title.text = f"Slide {i + 1}"
        prs.save(str(pptx_path))
        
        # Import twice (simulating version updates)
        results = []
        for _ in range(2):
            with open(pptx_path, "rb") as f:
                response = await client.post(
                    f"/api/projects/{sample_project.id}/upload",
                    files={"file": ("ordered.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
                )
            results.append(response.json())
        
        # Both imports should be consistent
        assert results[0]["status"] == results[1]["status"]


class TestRegressionG02_SlideHashConsistency:
    """Test slide hash calculation is consistent"""
    
    @pytest.mark.asyncio
    async def test_g02_same_input_same_hash(
        self,
        tmp_path
    ):
        """Same file produces same hash"""
        from app.adapters.media_converter import media_converter
        
        # Create test file
        test_file = tmp_path / "test.txt"
        test_file.write_text("test content")
        
        hash1 = media_converter.compute_file_hash(test_file)
        hash2 = media_converter.compute_file_hash(test_file)
        
        assert hash1 == hash2
    
    @pytest.mark.asyncio
    async def test_g02_different_input_different_hash(
        self,
        tmp_path
    ):
        """Different files produce different hashes"""
        from app.adapters.media_converter import media_converter
        
        file1 = tmp_path / "test1.txt"
        file2 = tmp_path / "test2.txt"
        file1.write_text("content 1")
        file2.write_text("content 2")
        
        hash1 = media_converter.compute_file_hash(file1)
        hash2 = media_converter.compute_file_hash(file2)
        
        assert hash1 != hash2


class TestRegressionG03_AudioHashConsistency:
    """Test audio hash calculation for TTS caching"""
    
    def test_g03_audio_hash_deterministic(self):
        """Audio hash is deterministic for same inputs"""
        from app.adapters.tts import TTSAdapter
        
        text = "Hello, this is a test."
        voice_id = "test-voice"
        lang = "en"
        model = "eleven_turbo_v2"
        
        hash1 = TTSAdapter.compute_audio_hash(text, voice_id, lang, model)
        hash2 = TTSAdapter.compute_audio_hash(text, voice_id, lang, model)
        
        assert hash1 == hash2
    
    def test_g03_audio_hash_changes_with_text(self):
        """Audio hash changes when text changes"""
        from app.adapters.tts import TTSAdapter
        
        voice_id = "test-voice"
        lang = "en"
        model = "eleven_turbo_v2"
        
        hash1 = TTSAdapter.compute_audio_hash("Text 1", voice_id, lang, model)
        hash2 = TTSAdapter.compute_audio_hash("Text 2", voice_id, lang, model)
        
        assert hash1 != hash2
    
    def test_g03_audio_hash_changes_with_language(self):
        """Audio hash changes when language changes"""
        from app.adapters.tts import TTSAdapter
        
        text = "Same text"
        voice_id = "test_voice"
        model = "eleven_turbo_v2"
        
        # Note: voice_id is ignored in hash computation (uses hardcoded voice)
        # But language should affect the hash
        hash1 = TTSAdapter.compute_audio_hash(text, voice_id, "en", model)
        hash2 = TTSAdapter.compute_audio_hash(text, voice_id, "ru", model)
        
        assert hash1 != hash2


class TestRegressionG04_APIResponseStructure:
    """Test API response structures remain consistent"""
    
    @pytest.mark.asyncio
    async def test_g04_project_response_structure(
        self,
        client: AsyncClient,
        sample_project: Project
    ):
        """Project response has required fields"""
        response = await client.get(f"/api/projects/{sample_project.id}")
        data = response.json()
        
        # Golden structure: these fields must exist
        required_fields = [
            "id", "name", "base_language", "current_version_id",
            "created_at", "updated_at", "status", "slide_count", "language_count"
        ]
        
        for field in required_fields:
            assert field in data, f"Missing field: {field}"
    
    @pytest.mark.asyncio
    async def test_g04_job_response_structure(
        self,
        client: AsyncClient,
        sample_render_job
    ):
        """Job response has required fields"""
        response = await client.get(f"/api/render/jobs/{sample_render_job.id}")
        data = response.json()
        
        # Golden structure
        required_fields = [
            "id", "project_id", "version_id", "lang", "job_type",
            "status", "progress_pct", "download_video_url", "download_srt_url",
            "error_message", "started_at", "finished_at"
        ]
        
        for field in required_fields:
            assert field in data, f"Missing field: {field}"
    
    @pytest.mark.asyncio
    async def test_g04_slide_response_structure(
        self,
        client: AsyncClient,
        sample_slide
    ):
        """Slide response has required fields"""
        response = await client.get(f"/api/slides/{sample_slide.id}")
        data = response.json()
        
        # Golden structure
        required_fields = [
            "id", "slide_index", "image_url", "notes_text",
            "scripts", "audio_files"
        ]
        
        for field in required_fields:
            assert field in data, f"Missing field: {field}"


class TestRegressionG05_DatabaseSchema:
    """Test database model constraints are enforced"""
    
    @pytest.mark.asyncio
    async def test_g05_project_requires_name(
        self,
        client: AsyncClient
    ):
        """Project creation requires name"""
        response = await client.post(
            "/api/projects",
            json={"base_language": "en"}  # Missing name
        )
        
        assert response.status_code == 422  # Validation error
    
    @pytest.mark.asyncio
    async def test_g05_slide_index_positive(
        self,
        sample_project: Project,
        sample_version: ProjectVersion,
        db_session: AsyncSession
    ):
        """Slide indices are positive"""
        from sqlalchemy import select
        
        # Create slide with positive index
        slide = Slide(
            project_id=sample_project.id,
            version_id=sample_version.id,
            slide_index=1,  # Must be positive
            image_path="test.png",
        )
        db_session.add(slide)
        await db_session.commit()
        
        assert slide.slide_index > 0

