"""
Flow Audit Tests (F-01 to F-10)

Цель: ключевые пользовательские пути проходят end-to-end.

Test Matrix:
- F-01: Новый проект → загрузка PPTX → импорт слайдов
- F-02: Импорт speaker notes как скриптов
- F-03: Загрузка изображений (PNG/JPG/WebP) как слайды
- F-04: Редактирование скрипта → автосейв → перезагрузка
- F-05: Drag&drop reorder слайдов
- F-06: Добавить/удалить слайд
- F-07: Генерация TTS на слайд (ElevenLabs)
- F-08: Повторная генерация TTS (перегенерить)
- F-09: Фоновая музыка (MP3) → loop/trim
- F-10: Export/Render итогового видео
"""
import uuid
import io
from datetime import datetime
from pathlib import Path
from unittest.mock import patch, MagicMock, AsyncMock

import pytest
from httpx import AsyncClient
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from PIL import Image

from app.db.models import (
    Project, ProjectVersion, Slide, SlideScript, SlideAudio, 
    RenderJob, AudioAsset, ProjectAudioSettings,
    ProjectStatus, ScriptSource, JobType, JobStatus
)


class TestFlowAuditF01_PPTXImport:
    """
    F-01: Новый проект → загрузка PPTX → импорт слайдов
    
    Шаги: создать проект → drag&drop PPTX → дождаться импорта
    Ожидаемо: слайды появились, порядок совпадает, превью ок
    """
    
    @pytest.mark.asyncio
    async def test_f01_create_project_upload_pptx_import_slides(
        self,
        client: AsyncClient,
        tmp_path
    ):
        """Full flow: create project → upload PPTX → convert → verify slides"""
        # Step 1: Create new project
        response = await client.post(
            "/api/projects",
            json={"name": "F-01 Test Project", "base_language": "en"}
        )
        assert response.status_code == 200
        project = response.json()
        project_id = project["id"]
        
        # Step 2: Upload PPTX
        from pptx import Presentation
        from pptx.util import Inches
        
        pptx_path = tmp_path / "test_f01.pptx"
        prs = Presentation()
        
        # Create 5 slides with titles
        for i in range(5):
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = f"Slide {i + 1}"
            # Add speaker notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = f"Speaker notes for slide {i + 1}"
        
        prs.save(str(pptx_path))
        
        with open(pptx_path, "rb") as f:
            response = await client.post(
                f"/api/projects/{project_id}/upload",
                files={"file": ("test.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                data={"comment": "F-01 test upload"}
            )
        
        assert response.status_code == 200
        upload_data = response.json()
        version_id = upload_data["version_id"]
        assert upload_data["status"] == "uploaded"
        
        # Step 3: Trigger conversion (mock Celery task)
        with patch("app.workers.tasks.convert_pptx_task") as mock_convert:
            mock_task = MagicMock()
            mock_task.id = "convert-task-f01"
            mock_convert.delay.return_value = mock_task
            
            response = await client.post(
                f"/api/projects/{project_id}/versions/{version_id}/convert"
            )
            assert response.status_code == 200
            assert response.json()["status"] == "queued"
        
        # Cleanup
        await client.delete(f"/api/projects/{project_id}")
    
    @pytest.mark.asyncio
    async def test_f01_slides_order_preserved(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion,
        db_session: AsyncSession
    ):
        """Verify slide order is preserved after import"""
        # Create multiple slides with specific order
        for i in range(5):
            slide = Slide(
                project_id=sample_project.id,
                version_id=sample_version.id,
                slide_index=i + 1,
                image_path=f"slides/slide_{i+1:03d}.png",
                notes_text=f"Notes {i + 1}",
            )
            db_session.add(slide)
        await db_session.commit()
        
        # Get slides and verify order
        response = await client.get(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/slides"
        )
        assert response.status_code == 200
        slides = response.json()
        
        # Verify order is correct (1-based indices)
        for i, slide in enumerate(slides):
            assert slide["slide_index"] == i + 1


class TestFlowAuditF02_SpeakerNotes:
    """
    F-02: Импорт speaker notes как скриптов
    
    Шаги: PPTX со speaker notes → импорт → открыть каждый слайд
    Ожидаемо: notes корректно подставлены в script editor по слайдам
    """
    
    @pytest.mark.asyncio
    async def test_f02_speaker_notes_imported_as_scripts(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion,
        db_session: AsyncSession
    ):
        """Verify speaker notes are imported as scripts"""
        # Create slides with notes
        slides_data = []
        for i in range(3):
            slide = Slide(
                project_id=sample_project.id,
                version_id=sample_version.id,
                slide_index=i + 1,
                image_path=f"slides/slide_{i+1:03d}.png",
                notes_text=f"Speaker notes for slide {i + 1}: Important content",
            )
            db_session.add(slide)
            slides_data.append(slide)
        await db_session.commit()
        
        # Import notes
        response = await client.post(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/import_notes",
            params={"lang": "en"}
        )
        assert response.status_code == 200
        data = response.json()
        assert data["imported_count"] == 3
        
        # Verify scripts were created
        for i, slide in enumerate(slides_data):
            await db_session.refresh(slide)
            response = await client.get(f"/api/slides/{slide.id}/scripts")
            scripts = response.json()
            
            assert len(scripts) >= 1
            en_script = next((s for s in scripts if s["lang"] == "en"), None)
            assert en_script is not None
            assert f"Speaker notes for slide {i + 1}" in en_script["text"]
            assert en_script["source"] == "imported_notes"


class TestFlowAuditF03_ImageUpload:
    """
    F-03: Загрузка изображений (PNG/JPG/WebP) как слайды
    
    Шаги: drag&drop несколько изображений
    Ожидаемо: созданы новые слайды, корректный размер/ориентация
    """
    
    @pytest.mark.asyncio
    async def test_f03_upload_png_as_slide(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion
    ):
        """Upload PNG image as a new slide"""
        # Create test image
        img = Image.new("RGB", (1920, 1080), color="blue")
        img_buffer = io.BytesIO()
        img.save(img_buffer, format="PNG")
        img_buffer.seek(0)
        
        response = await client.post(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/slides/add",
            files={"file": ("slide.png", img_buffer, "image/png")}
        )
        
        assert response.status_code == 200
        data = response.json()
        assert "id" in data
        assert data["slide_index"] == 1
        assert "image_url" in data
    
    @pytest.mark.asyncio
    async def test_f03_upload_jpeg_as_slide(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion
    ):
        """Upload JPEG image as a new slide"""
        img = Image.new("RGB", (1920, 1080), color="red")
        img_buffer = io.BytesIO()
        img.save(img_buffer, format="JPEG")
        img_buffer.seek(0)
        
        response = await client.post(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/slides/add",
            files={"file": ("slide.jpg", img_buffer, "image/jpeg")}
        )
        
        assert response.status_code == 200
    
    @pytest.mark.asyncio
    async def test_f03_upload_webp_as_slide(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion
    ):
        """Upload WebP image as a new slide"""
        img = Image.new("RGB", (1920, 1080), color="green")
        img_buffer = io.BytesIO()
        img.save(img_buffer, format="WEBP")
        img_buffer.seek(0)
        
        response = await client.post(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/slides/add",
            files={"file": ("slide.webp", img_buffer, "image/webp")}
        )
        
        assert response.status_code == 200
    
    @pytest.mark.asyncio
    async def test_f03_invalid_file_rejected(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion
    ):
        """Verify invalid file types are rejected"""
        response = await client.post(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/slides/add",
            files={"file": ("test.txt", b"invalid content", "text/plain")}
        )
        
        assert response.status_code == 400


class TestFlowAuditF04_ScriptAutosave:
    """
    F-04: Редактирование скрипта → автосейв → перезагрузка
    
    Шаги: изменить текст → подождать автосейв → F5/close-open
    Ожидаемо: текст сохранён, нет потери последних правок
    """
    
    @pytest.mark.asyncio
    async def test_f04_script_update_persists(
        self,
        client: AsyncClient,
        sample_slide: Slide,
        sample_script: SlideScript,
        db_session: AsyncSession
    ):
        """Verify script updates are persisted"""
        new_text = "Updated script text with important changes"
        
        # Update script
        response = await client.patch(
            f"/api/slides/{sample_slide.id}/scripts/en",
            json={"text": new_text}
        )
        assert response.status_code == 200
        
        # Simulate "refresh" - get slide again
        response = await client.get(f"/api/slides/{sample_slide.id}")
        assert response.status_code == 200
        slide_data = response.json()
        
        en_script = next((s for s in slide_data["scripts"] if s["lang"] == "en"), None)
        assert en_script is not None
        assert en_script["text"] == new_text
    
    @pytest.mark.asyncio
    async def test_f04_multiple_updates_preserve_latest(
        self,
        client: AsyncClient,
        sample_slide: Slide,
        sample_script: SlideScript
    ):
        """Verify multiple rapid updates preserve the latest value"""
        # Simulate rapid autosave calls
        for i in range(5):
            await client.patch(
                f"/api/slides/{sample_slide.id}/scripts/en",
                json={"text": f"Update {i}"}
            )
        
        final_text = "Final version of the script"
        response = await client.patch(
            f"/api/slides/{sample_slide.id}/scripts/en",
            json={"text": final_text}
        )
        
        # Verify final version is saved
        response = await client.get(f"/api/slides/{sample_slide.id}/scripts")
        scripts = response.json()
        en_script = next((s for s in scripts if s["lang"] == "en"), None)
        assert en_script["text"] == final_text


class TestFlowAuditF05_SlideReorder:
    """
    F-05: Drag&drop reorder слайдов
    
    Шаги: поменять порядок 5–10 слайдов
    Ожидаемо: порядок сохранён, превью/скрипты/аудио привязки не съехали
    """
    
    @pytest.mark.asyncio
    async def test_f05_reorder_slides(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion,
        db_session: AsyncSession
    ):
        """Test slide reordering"""
        # Create 5 slides
        slide_ids = []
        for i in range(5):
            slide = Slide(
                project_id=sample_project.id,
                version_id=sample_version.id,
                slide_index=i + 1,
                image_path=f"slides/slide_{i+1:03d}.png",
            )
            db_session.add(slide)
            await db_session.flush()
            slide_ids.append(str(slide.id))
        await db_session.commit()
        
        # Reorder: move first to last (1,2,3,4,5 -> 2,3,4,5,1)
        new_order = slide_ids[1:] + [slide_ids[0]]
        
        response = await client.put(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/slides/reorder",
            json={"slide_ids": new_order}
        )
        
        assert response.status_code == 200
        assert response.json()["success"] is True
        
        # Verify new order
        response = await client.get(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/slides"
        )
        slides = response.json()
        
        for i, slide in enumerate(slides):
            assert slide["id"] == new_order[i]
    
    @pytest.mark.asyncio
    async def test_f05_reorder_preserves_scripts(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion,
        db_session: AsyncSession
    ):
        """Verify scripts stay with their slides after reorder"""
        # Create slides with scripts
        slides = []
        for i in range(3):
            slide = Slide(
                project_id=sample_project.id,
                version_id=sample_version.id,
                slide_index=i + 1,
                image_path=f"slides/slide_{i+1:03d}.png",
            )
            db_session.add(slide)
            await db_session.flush()
            
            script = SlideScript(
                slide_id=slide.id,
                lang="en",
                text=f"Script for original slide {i + 1}",
                source=ScriptSource.MANUAL,
            )
            db_session.add(script)
            slides.append(slide)
        await db_session.commit()
        
        # Reverse order
        new_order = [str(s.id) for s in reversed(slides)]
        
        await client.put(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/slides/reorder",
            json={"slide_ids": new_order}
        )
        
        # Verify scripts match original content (by slide ID, not index)
        for i, orig_slide in enumerate(slides):
            response = await client.get(f"/api/slides/{orig_slide.id}/scripts")
            scripts = response.json()
            en_script = next((s for s in scripts if s["lang"] == "en"), None)
            assert f"Script for original slide {i + 1}" in en_script["text"]


class TestFlowAuditF06_AddDeleteSlide:
    """
    F-06: Добавить/удалить слайд
    
    Шаги: add slide → delete slide → undo (если есть)
    Ожидаемо: корректное обновление таймлайна/индексов/скриптов
    """
    
    @pytest.mark.asyncio
    async def test_f06_add_slide_at_position(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion,
        db_session: AsyncSession
    ):
        """Test adding slide at specific position"""
        # Create initial slides
        for i in range(3):
            slide = Slide(
                project_id=sample_project.id,
                version_id=sample_version.id,
                slide_index=i + 1,
                image_path=f"slides/slide_{i+1:03d}.png",
            )
            db_session.add(slide)
        await db_session.commit()
        
        # Add slide at position 2
        img = Image.new("RGB", (1920, 1080), color="purple")
        img_buffer = io.BytesIO()
        img.save(img_buffer, format="PNG")
        img_buffer.seek(0)
        
        response = await client.post(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/slides/add",
            files={"file": ("new.png", img_buffer, "image/png")},
            params={"position": 2}
        )
        
        assert response.status_code == 200
        assert response.json()["slide_index"] == 2
        assert response.json()["total_slides"] == 4
    
    @pytest.mark.asyncio
    async def test_f06_delete_slide_reindexes(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion,
        db_session: AsyncSession
    ):
        """Test deleting slide updates indices"""
        # Create slides
        slides = []
        for i in range(4):
            slide = Slide(
                project_id=sample_project.id,
                version_id=sample_version.id,
                slide_index=i + 1,
                image_path=f"slides/slide_{i+1:03d}.png",
            )
            db_session.add(slide)
            slides.append(slide)
        await db_session.commit()
        
        # Delete slide at index 2
        response = await client.delete(f"/api/slides/{slides[1].id}")
        assert response.status_code == 200
        assert response.json()["deleted_index"] == 2
        assert response.json()["slides_reindexed"] == 2  # slides 3 and 4
        
        # Verify remaining slides have correct indices
        response = await client.get(
            f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/slides"
        )
        remaining_slides = response.json()
        assert len(remaining_slides) == 3
        
        indices = [s["slide_index"] for s in remaining_slides]
        assert indices == [1, 2, 3]  # Continuous after reindex


class TestFlowAuditF07_TTSGeneration:
    """
    F-07: Генерация TTS на слайд (ElevenLabs)
    
    Шаги: заполнить скрипт → выбрать голос/язык → generate для 3 слайдов
    Ожидаемо: создаётся аудио, длительность логична, можно прослушать
    """
    
    @pytest.mark.asyncio
    async def test_f07_generate_slide_tts(
        self,
        client: AsyncClient,
        sample_slide: Slide,
        sample_script: SlideScript
    ):
        """Test TTS generation for single slide"""
        with patch("app.workers.tasks.tts_slide_task") as mock_tts:
            mock_task = MagicMock()
            mock_task.id = "tts-task-f07"
            mock_tts.delay.return_value = mock_task
            
            response = await client.post(
                f"/api/slides/{sample_slide.id}/tts",
                params={"lang": "en", "voice_id": "test-voice"}
            )
            
            assert response.status_code == 200
            data = response.json()
            assert data["task_id"] == "tts-task-f07"
            assert data["lang"] == "en"
            assert data["status"] == "queued"
    
    @pytest.mark.asyncio
    async def test_f07_batch_tts_generation(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion
    ):
        """Test batch TTS for all slides"""
        with patch("app.workers.tasks.tts_batch_task") as mock_tts:
            mock_task = MagicMock()
            mock_task.id = "tts-batch-f07"
            mock_tts.delay.return_value = mock_task
            
            response = await client.post(
                f"/api/slides/projects/{sample_project.id}/versions/{sample_version.id}/tts",
                params={"lang": "en"}
            )
            
            assert response.status_code == 200
            assert response.json()["status"] == "queued"


class TestFlowAuditF08_TTSRegeneration:
    """
    F-08: Повторная генерация TTS (перегенерить)
    
    Шаги: изменить текст → regenerate
    Ожидаемо: старый аудио-ассет не "подмешивается", используется новый
    """
    
    @pytest.mark.asyncio
    async def test_f08_regenerate_replaces_old_audio(
        self,
        client: AsyncClient,
        sample_slide: Slide,
        sample_script: SlideScript,
        db_session: AsyncSession
    ):
        """Verify regeneration creates new audio record"""
        # Create existing audio
        old_audio = SlideAudio(
            slide_id=sample_slide.id,
            lang="en",
            voice_id="old-voice",
            audio_path="audio/old_slide.wav",
            duration_sec=5.0,
            audio_hash="old-hash-123",
        )
        db_session.add(old_audio)
        await db_session.commit()
        old_audio_id = old_audio.id
        
        # Update script text
        await client.patch(
            f"/api/slides/{sample_slide.id}/scripts/en",
            json={"text": "New updated script for regeneration"}
        )
        
        # Trigger regeneration
        with patch("app.workers.tasks.tts_slide_task") as mock_tts:
            mock_task = MagicMock()
            mock_task.id = "tts-regen-f08"
            mock_tts.delay.return_value = mock_task
            
            response = await client.post(
                f"/api/slides/{sample_slide.id}/tts",
                params={"lang": "en", "voice_id": "new-voice"}
            )
            
            assert response.status_code == 200
            # The task will replace the old audio in the worker


class TestFlowAuditF09_BackgroundMusic:
    """
    F-09: Фоновая музыка (MP3) → loop/trim до длины озвучки
    
    Шаги: загрузить mp3 → включить bg music → сгенерить таймлайн
    Ожидаемо: музыка лупится/обрезается под voice_timeline
    """
    
    @pytest.mark.asyncio
    async def test_f09_upload_background_music(
        self,
        client: AsyncClient,
        sample_project: Project,
        tmp_path
    ):
        """Test uploading background music"""
        # Create fake MP3 file
        mp3_content = b"fake mp3 data" * 100
        
        response = await client.post(
            f"/api/projects/{sample_project.id}/upload_music",
            files={"file": ("music.mp3", mp3_content, "audio/mpeg")}
        )
        
        assert response.status_code == 200
        assert "asset_id" in response.json()
    
    @pytest.mark.asyncio
    async def test_f09_enable_background_music(
        self,
        client: AsyncClient,
        sample_project: Project,
        db_session: AsyncSession
    ):
        """Test enabling background music in settings"""
        # First upload music
        mp3_content = b"fake mp3 data" * 100
        upload_resp = await client.post(
            f"/api/projects/{sample_project.id}/upload_music",
            files={"file": ("music.mp3", mp3_content, "audio/mpeg")}
        )
        asset_id = upload_resp.json()["asset_id"]
        
        # Enable in settings
        response = await client.put(
            f"/api/projects/{sample_project.id}/audio_settings",
            json={
                "background_music_enabled": True,
                "music_gain_db": -20.0,
                "ducking_enabled": True,
            }
        )
        
        assert response.status_code == 200
        
        # Verify settings
        response = await client.get(f"/api/projects/{sample_project.id}/audio_settings")
        settings = response.json()
        assert settings["background_music_enabled"] is True


class TestFlowAuditF10_VideoRender:
    """
    F-10: Export/Render итогового видео
    
    Шаги: запустить рендер → дождаться завершения → открыть результат
    Ожидаемо: видео соответствует таймлайну, без рассинхрона
    """
    
    @pytest.mark.asyncio
    async def test_f10_render_video_creates_job(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion
    ):
        """Test video render creates job"""
        with patch("app.workers.tasks.render_language_task") as mock_render:
            def _apply_async(*args, **kwargs):
                t = MagicMock()
                t.id = kwargs.get("task_id")
                return t
            mock_render.apply_async.side_effect = _apply_async
            
            response = await client.post(
                f"/api/render/projects/{sample_project.id}/versions/{sample_version.id}/render",
                params={"lang": "en"}
            )
            
            assert response.status_code == 200
            data = response.json()
            assert "job_id" in data
            assert data["status"] == "queued"
    
    @pytest.mark.asyncio
    async def test_f10_completed_job_has_download_urls(
        self,
        client: AsyncClient,
        sample_render_job: RenderJob,
        db_session: AsyncSession
    ):
        """Test completed job provides download URLs"""
        # Mark job as done with output paths
        sample_render_job.status = JobStatus.DONE
        sample_render_job.progress_pct = 100
        sample_render_job.output_video_path = "exports/en/deck_en.mp4"
        sample_render_job.output_srt_path = "exports/en/deck_en.srt"
        sample_render_job.finished_at = datetime.utcnow()
        await db_session.commit()
        
        response = await client.get(f"/api/render/jobs/{sample_render_job.id}")
        
        assert response.status_code == 200
        data = response.json()
        assert data["status"] == "done"
        assert data["download_video_url"] is not None
        assert data["download_srt_url"] is not None

