"""
Tests for PPTXConverter
"""
import hashlib
import pytest
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

from app.adapters.pptx_converter import PPTXConverter


class TestPPTXConverter:
    """Tests for PPTXConverter class"""
    
    @pytest.fixture
    def converter(self):
        """Create converter instance"""
        return PPTXConverter(dpi=150)
    
    def test_init_default_dpi(self):
        """Test default DPI setting"""
        converter = PPTXConverter()
        assert converter.dpi == 150
    
    def test_init_custom_dpi(self):
        """Test custom DPI setting"""
        converter = PPTXConverter(dpi=300)
        assert converter.dpi == 300
    
    def test_compute_slide_hash(self, converter, tmp_path):
        """Test slide image hash computation"""
        # Create test image file
        image_path = tmp_path / "slide.png"
        image_path.write_bytes(b"fake image data")
        
        hash1 = converter.compute_slide_hash(image_path)
        hash2 = converter.compute_slide_hash(image_path)
        
        assert hash1 == hash2
        assert len(hash1) == 64  # SHA256 hex length
    
    def test_compute_slide_hash_different_content(self, converter, tmp_path):
        """Test that different images produce different hashes"""
        image1 = tmp_path / "slide1.png"
        image1.write_bytes(b"image data 1")
        
        image2 = tmp_path / "slide2.png"
        image2.write_bytes(b"image data 2")
        
        hash1 = converter.compute_slide_hash(image1)
        hash2 = converter.compute_slide_hash(image2)
        
        assert hash1 != hash2


class TestPPTXConverterWithRealFiles:
    """Tests that require creating real PPTX files"""
    
    @pytest.fixture
    def converter(self):
        return PPTXConverter(dpi=150)
    
    @pytest.fixture
    def test_pptx(self, tmp_path):
        """Create a test PPTX file"""
        from pptx import Presentation
        
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        
        # Add 3 slides with notes
        for i in range(3):
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            if title:
                title.text = f"Slide {i + 1}"
            
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = f"Notes for slide {i + 1}"
        
        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))
        return pptx_path
    
    def test_extract_speaker_notes(self, converter, test_pptx):
        """Test speaker notes extraction"""
        notes = converter.extract_speaker_notes(test_pptx)
        
        assert len(notes) == 3
        assert notes[0] == "Notes for slide 1"
        assert notes[1] == "Notes for slide 2"
        assert notes[2] == "Notes for slide 3"
    
    def test_extract_speaker_notes_empty(self, converter, tmp_path):
        """Test extraction when slides have no notes"""
        from pptx import Presentation
        
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        
        # Add slide without notes
        prs.slides.add_slide(slide_layout)
        
        pptx_path = tmp_path / "no_notes.pptx"
        prs.save(str(pptx_path))
        
        notes = converter.extract_speaker_notes(pptx_path)
        
        assert len(notes) == 1
        # Empty notes are returned as None
        assert notes[0] is None or notes[0] == ""
    
    def test_compute_slides_hash(self, converter, test_pptx):
        """Test slides hash computation"""
        hash1 = converter.compute_slides_hash(test_pptx)
        hash2 = converter.compute_slides_hash(test_pptx)
        
        assert hash1 == hash2
        assert len(hash1) == 64
    
    def test_compute_slides_hash_changes_with_content(self, converter, tmp_path):
        """Test that hash changes when PPTX structure changes"""
        from pptx import Presentation
        
        # Create first PPTX with 2 slides
        prs1 = Presentation()
        prs1.slides.add_slide(prs1.slide_layouts[0])
        prs1.slides.add_slide(prs1.slide_layouts[0])
        
        pptx1 = tmp_path / "pptx1.pptx"
        prs1.save(str(pptx1))
        
        # Create second PPTX with 3 slides
        prs2 = Presentation()
        prs2.slides.add_slide(prs2.slide_layouts[0])
        prs2.slides.add_slide(prs2.slide_layouts[0])
        prs2.slides.add_slide(prs2.slide_layouts[0])
        
        pptx2 = tmp_path / "pptx2.pptx"
        prs2.save(str(pptx2))
        
        hash1 = converter.compute_slides_hash(pptx1)
        hash2 = converter.compute_slides_hash(pptx2)
        
        assert hash1 != hash2


class TestPPTXConverterConversion:
    """Tests for PPTX to PNG conversion"""
    
    @pytest.fixture
    def converter(self):
        return PPTXConverter(dpi=150)
    
    @pytest.mark.asyncio
    async def test_convert_to_pdf(self, converter, tmp_path):
        """Test PPTX to PDF conversion command"""
        pptx_path = tmp_path / "test.pptx"
        output_dir = tmp_path / "output"
        output_dir.mkdir()
        
        # Mock subprocess
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b'', b''))
        mock_process.returncode = 0
        
        with patch('asyncio.create_subprocess_exec', return_value=mock_process) as mock_exec:
            await converter._convert_to_pdf(pptx_path, output_dir)
            
            # Verify correct command was called
            mock_exec.assert_called_once()
            call_args = mock_exec.call_args[0]
            
            assert "libreoffice" in call_args
            assert "--headless" in call_args
            assert "--convert-to" in call_args
            assert "pdf" in call_args
    
    @pytest.mark.asyncio
    async def test_convert_to_pdf_failure(self, converter, tmp_path):
        """Test PDF conversion failure handling"""
        pptx_path = tmp_path / "test.pptx"
        output_dir = tmp_path / "output"
        
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b'', b'Error message'))
        mock_process.returncode = 1
        
        with patch('asyncio.create_subprocess_exec', return_value=mock_process):
            with pytest.raises(RuntimeError, match="LibreOffice conversion failed"):
                await converter._convert_to_pdf(pptx_path, output_dir)
    
    @pytest.mark.asyncio
    async def test_pdf_to_png_pdftoppm(self, converter, tmp_path):
        """Test PDF to PNG conversion using pdftoppm"""
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(b"fake pdf")
        output_dir = tmp_path / "output"
        output_dir.mkdir()
        
        # Create mock output files
        (output_dir / "slide-1.png").write_bytes(b"png1")
        (output_dir / "slide-2.png").write_bytes(b"png2")
        
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b'', b''))
        mock_process.returncode = 0
        
        with patch('asyncio.create_subprocess_exec', return_value=mock_process) as mock_exec:
            result = await converter._pdf_to_png_pdftoppm(pdf_path, output_dir)
            
            # Verify pdftoppm was called
            mock_exec.assert_called_once()
            call_args = mock_exec.call_args[0]
            assert "pdftoppm" in call_args
            assert "-png" in call_args
            
            # Verify files were renamed
            assert len(result) == 2
            assert result[0].name == "001.png"
            assert result[1].name == "002.png"
    
    @pytest.mark.asyncio
    async def test_pdf_to_png_imagemagick_fallback(self, converter, tmp_path):
        """Test fallback to ImageMagick when pdftoppm not available"""
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(b"fake pdf")
        output_dir = tmp_path / "output"
        output_dir.mkdir()
        
        # Create mock output file
        (output_dir / "000.png").write_bytes(b"png")
        
        mock_process = AsyncMock()
        mock_process.communicate = AsyncMock(return_value=(b'', b''))
        mock_process.returncode = 0
        
        # Make pdftoppm fail with FileNotFoundError
        async def mock_pdftoppm_fails(*args, **kwargs):
            raise FileNotFoundError("pdftoppm not found")
        
        with patch.object(
            converter,
            '_pdf_to_png_pdftoppm',
            side_effect=FileNotFoundError
        ):
            with patch('asyncio.create_subprocess_exec', return_value=mock_process) as mock_exec:
                result = await converter._pdf_to_png(pdf_path, output_dir)
                
                # Verify ImageMagick convert was called
                mock_exec.assert_called()
    
    @pytest.mark.asyncio
    async def test_convert_to_png_full_pipeline(self, converter, tmp_path):
        """Test full conversion pipeline"""
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(b"fake pptx")
        output_dir = tmp_path / "slides"
        
        # Mock the entire pipeline
        async def mock_convert_to_pdf(pptx, out_dir):
            (out_dir / "presentation.pdf").write_bytes(b"pdf")
        
        async def mock_pdf_to_png(pdf, out_dir):
            png1 = out_dir / "001.png"
            png2 = out_dir / "002.png"
            png1.write_bytes(b"png1")
            png2.write_bytes(b"png2")
            return [png1, png2]
        
        with patch.object(converter, '_convert_to_pdf', mock_convert_to_pdf):
            with patch.object(converter, '_pdf_to_png', mock_pdf_to_png):
                result = await converter.convert_to_png(pptx_path, output_dir)
                
                assert len(result) == 2
                assert all(p.suffix == ".png" for p in result)


class TestPPTXConverterEdgeCases:
    """Edge case tests for PPTXConverter"""
    
    @pytest.fixture
    def converter(self):
        return PPTXConverter()
    
    def test_extract_notes_mixed(self, converter, tmp_path):
        """Test extraction when only some slides have notes"""
        from pptx import Presentation
        
        prs = Presentation()
        
        # Slide 1 with notes
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.notes_slide.notes_text_frame.text = "Notes for slide 1"
        
        # Slide 2 without notes (just add slide, don't set notes)
        slide2 = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Slide 3 with notes
        slide3 = prs.slides.add_slide(prs.slide_layouts[0])
        slide3.notes_slide.notes_text_frame.text = "Notes for slide 3"
        
        pptx_path = tmp_path / "mixed.pptx"
        prs.save(str(pptx_path))
        
        notes = converter.extract_speaker_notes(pptx_path)
        
        assert len(notes) == 3
        assert notes[0] == "Notes for slide 1"
        assert notes[1] is None or notes[1] == ""
        assert notes[2] == "Notes for slide 3"
    
    def test_extract_notes_multiline(self, converter, tmp_path):
        """Test extraction of multiline notes"""
        from pptx import Presentation
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Line 1\nLine 2\nLine 3"
        
        pptx_path = tmp_path / "multiline.pptx"
        prs.save(str(pptx_path))
        
        notes = converter.extract_speaker_notes(pptx_path)
        
        assert len(notes) == 1
        assert "Line 1" in notes[0]
        assert "Line 2" in notes[0]
        assert "Line 3" in notes[0]

