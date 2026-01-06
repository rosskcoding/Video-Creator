"""
Test configuration and fixtures
"""
import asyncio
import base64
import os
import uuid
from datetime import datetime
from pathlib import Path
from typing import AsyncGenerator, Generator
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
import pytest_asyncio
from fastapi.testclient import TestClient
from httpx import AsyncClient, ASGITransport
from sqlalchemy import create_engine, text
from sqlalchemy.ext.asyncio import AsyncSession, create_async_engine, async_sessionmaker
from sqlalchemy.orm import sessionmaker
from sqlalchemy.pool import StaticPool

# Set test environment
os.environ["DATABASE_URL"] = "sqlite+aiosqlite:///:memory:"
os.environ["DATA_DIR"] = "/tmp/test_data"
os.environ["OPENAI_API_KEY"] = "test-key"
os.environ["ELEVENLABS_API_KEY"] = "test-key"
os.environ["DEBUG"] = "true"
os.environ["ADMIN_USERNAME"] = "login"
os.environ["ADMIN_PASSWORD"] = "Superman2026!"


def _basic_auth_header(username: str, password: str) -> str:
    token = base64.b64encode(f"{username}:{password}".encode("utf-8")).decode("utf-8")
    return f"Basic {token}"


AUTH_HEADERS = {
    "Authorization": _basic_auth_header(
        os.environ["ADMIN_USERNAME"],
        os.environ["ADMIN_PASSWORD"],
    )
}

from app.db.database import Base, get_db
from app.main import app
from app.db.models import (
    Project, ProjectVersion, ProjectAudioSettings, ProjectTranslationRules,
    Slide, SlideScript, SlideAudio, AudioAsset, RenderJob,
    ProjectStatus, ScriptSource, JobType, JobStatus, DuckingStrength, TranslationStyle
)


# Test database setup
TEST_DATABASE_URL = "sqlite+aiosqlite:///:memory:"


@pytest.fixture(scope="session")
def event_loop():
    """Create event loop for async tests"""
    loop = asyncio.new_event_loop()
    yield loop
    loop.close()


@pytest_asyncio.fixture(scope="function")
async def async_engine():
    """Create async engine for tests"""
    engine = create_async_engine(
        TEST_DATABASE_URL,
        connect_args={"check_same_thread": False},
        poolclass=StaticPool,
        echo=False,
    )
    
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
    
    yield engine
    
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.drop_all)
    
    await engine.dispose()


@pytest_asyncio.fixture(scope="function")
async def db_session(async_engine) -> AsyncGenerator[AsyncSession, None]:
    """Create database session for tests"""
    async_session = async_sessionmaker(
        async_engine,
        class_=AsyncSession,
        expire_on_commit=False,
        autocommit=False,
        autoflush=False,
    )
    
    async with async_session() as session:
        yield session
        await session.rollback()


@pytest_asyncio.fixture(scope="function")
async def client(db_session: AsyncSession) -> AsyncGenerator[AsyncClient, None]:
    """Create test client with database session override"""
    
    async def override_get_db():
        yield db_session
    
    app.dependency_overrides[get_db] = override_get_db
    
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        ac.headers.update(AUTH_HEADERS)
        yield ac
    
    app.dependency_overrides.clear()


@pytest.fixture
def sync_client(db_session: AsyncSession) -> Generator[TestClient, None, None]:
    """Create synchronous test client"""
    
    async def override_get_db():
        yield db_session
    
    app.dependency_overrides[get_db] = override_get_db
    
    with TestClient(app) as tc:
        tc.headers.update(AUTH_HEADERS)
        yield tc
    
    app.dependency_overrides.clear()


# === Sample Data Fixtures ===

@pytest_asyncio.fixture
async def sample_project(db_session: AsyncSession) -> Project:
    """Create a sample project"""
    project = Project(
        name="Test Project",
        base_language="en",
        allowed_languages=["en"],
    )
    db_session.add(project)
    await db_session.flush()  # Flush to get the project ID
    
    # Create audio settings
    audio_settings = ProjectAudioSettings(project_id=project.id)
    db_session.add(audio_settings)
    
    # Create translation rules
    translation_rules = ProjectTranslationRules(project_id=project.id)
    db_session.add(translation_rules)
    
    await db_session.commit()
    await db_session.refresh(project)
    return project


@pytest_asyncio.fixture
async def sample_version(db_session: AsyncSession, sample_project: Project) -> ProjectVersion:
    """Create a sample project version"""
    version = ProjectVersion(
        project_id=sample_project.id,
        version_number=1,
        status=ProjectStatus.DRAFT,
        comment="Initial version",
    )
    db_session.add(version)
    
    # Update project current version
    sample_project.current_version_id = version.id
    
    await db_session.commit()
    await db_session.refresh(version)
    return version


@pytest_asyncio.fixture
async def sample_slide(
    db_session: AsyncSession,
    sample_project: Project,
    sample_version: ProjectVersion
) -> Slide:
    """Create a sample slide"""
    slide = Slide(
        project_id=sample_project.id,
        version_id=sample_version.id,
        slide_index=1,
        image_path="/tmp/test_data/slide_001.png",
        notes_text="Speaker notes for slide 1",
    )
    db_session.add(slide)
    await db_session.commit()
    await db_session.refresh(slide)
    return slide


@pytest_asyncio.fixture
async def sample_script(db_session: AsyncSession, sample_slide: Slide) -> SlideScript:
    """Create a sample script"""
    script = SlideScript(
        slide_id=sample_slide.id,
        lang="en",
        text="This is the script text for slide 1",
        source=ScriptSource.MANUAL,
    )
    db_session.add(script)
    await db_session.commit()
    await db_session.refresh(script)
    return script


@pytest_asyncio.fixture
async def sample_render_job(
    db_session: AsyncSession,
    sample_project: Project,
    sample_version: ProjectVersion
) -> RenderJob:
    """Create a sample render job"""
    job = RenderJob(
        project_id=sample_project.id,
        version_id=sample_version.id,
        lang="en",
        job_type=JobType.RENDER,
        status=JobStatus.QUEUED,
    )
    db_session.add(job)
    await db_session.commit()
    await db_session.refresh(job)
    return job


# === Mock Fixtures ===

@pytest.fixture
def mock_celery_task():
    """Mock Celery task for testing async task calls"""
    mock_task = MagicMock()
    mock_task.id = "test-task-id-123"
    mock_delay = MagicMock(return_value=mock_task)
    return mock_delay


@pytest.fixture
def mock_openai_client():
    """Mock OpenAI client for translation tests"""
    mock_response = MagicMock()
    mock_response.choices = [MagicMock()]
    mock_response.choices[0].message.content = "Translated text"
    
    mock_client = AsyncMock()
    mock_client.chat.completions.create = AsyncMock(return_value=mock_response)
    return mock_client


@pytest.fixture
def mock_elevenlabs_client():
    """Mock ElevenLabs client for TTS tests"""
    mock_client = MagicMock()
    mock_client.text_to_speech.convert = MagicMock(
        return_value=[b"audio", b"data", b"chunks"]
    )
    return mock_client


@pytest.fixture
def temp_dir(tmp_path):
    """Create a temporary directory for file operations"""
    return tmp_path


# === Helper Functions ===

def create_test_pptx(path: Path) -> None:
    """Create a minimal test PPTX file"""
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    
    # Add a few slides
    for i in range(3):
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title:
            title.text = f"Slide {i + 1}"
        
        # Add notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"Notes for slide {i + 1}"
    
    prs.save(str(path))

