"""
Tests for Project management API routes
"""
import uuid
import pytest
from unittest.mock import patch, MagicMock, AsyncMock

from httpx import AsyncClient
from sqlalchemy.ext.asyncio import AsyncSession

from app.db.models import Project, ProjectVersion, ProjectAudioSettings, ProjectTranslationRules


class TestProjectCRUD:
    """Tests for project CRUD operations"""
    
    @pytest.mark.asyncio
    async def test_create_project(self, client: AsyncClient):
        """Test creating a new project"""
        response = await client.post(
            "/api/projects",
            json={"name": "Test Project", "base_language": "en"}
        )
        
        assert response.status_code == 200
        data = response.json()
        
        assert data["name"] == "Test Project"
        assert data["base_language"] == "en"
        assert "id" in data
        assert "created_at" in data
        assert "updated_at" in data
    
    @pytest.mark.asyncio
    async def test_create_project_default_language(self, client: AsyncClient):
        """Test creating project with default base_language"""
        response = await client.post(
            "/api/projects",
            json={"name": "Default Lang Project"}
        )
        
        assert response.status_code == 200
        data = response.json()
        
        assert data["base_language"] == "en"
    
    @pytest.mark.asyncio
    async def test_list_projects_empty(self, client: AsyncClient):
        """Test listing projects when empty"""
        response = await client.get("/api/projects")
        
        assert response.status_code == 200
        assert response.json() == []
    
    @pytest.mark.asyncio
    async def test_list_projects(self, client: AsyncClient, sample_project: Project):
        """Test listing projects"""
        response = await client.get("/api/projects")
        
        assert response.status_code == 200
        data = response.json()
        
        assert len(data) == 1
        assert data[0]["name"] == "Test Project"
    
    @pytest.mark.asyncio
    async def test_get_project(self, client: AsyncClient, sample_project: Project):
        """Test getting a single project"""
        response = await client.get(f"/api/projects/{sample_project.id}")
        
        assert response.status_code == 200
        data = response.json()
        
        assert data["id"] == str(sample_project.id)
        assert data["name"] == sample_project.name
    
    @pytest.mark.asyncio
    async def test_get_project_not_found(self, client: AsyncClient):
        """Test getting a non-existent project"""
        fake_id = uuid.uuid4()
        response = await client.get(f"/api/projects/{fake_id}")
        
        assert response.status_code == 404
        assert "not found" in response.json()["detail"].lower()
    
    @pytest.mark.asyncio
    async def test_update_project(self, client: AsyncClient, sample_project: Project):
        """Test updating a project"""
        response = await client.patch(
            f"/api/projects/{sample_project.id}",
            json={"name": "Updated Name", "base_language": "ru"}
        )
        
        assert response.status_code == 200
        data = response.json()
        
        assert data["name"] == "Updated Name"
        assert data["base_language"] == "ru"
    
    @pytest.mark.asyncio
    async def test_update_project_partial(self, client: AsyncClient, sample_project: Project):
        """Test partial project update"""
        response = await client.patch(
            f"/api/projects/{sample_project.id}",
            json={"name": "Only Name Changed"}
        )
        
        assert response.status_code == 200
        data = response.json()
        
        assert data["name"] == "Only Name Changed"
        assert data["base_language"] == "en"  # Should remain unchanged
    
    @pytest.mark.asyncio
    async def test_delete_project(self, client: AsyncClient, sample_project: Project):
        """Test deleting a project"""
        response = await client.delete(f"/api/projects/{sample_project.id}")
        
        assert response.status_code == 200
        assert response.json()["status"] == "deleted"
        
        # Verify project is gone
        response = await client.get(f"/api/projects/{sample_project.id}")
        assert response.status_code == 404
    
    @pytest.mark.asyncio
    async def test_delete_project_not_found(self, client: AsyncClient):
        """Test deleting a non-existent project"""
        fake_id = uuid.uuid4()
        response = await client.delete(f"/api/projects/{fake_id}")
        
        assert response.status_code == 404


class TestVersions:
    """Tests for project version management"""
    
    @pytest.mark.asyncio
    async def test_list_versions_empty(
        self,
        client: AsyncClient,
        sample_project: Project
    ):
        """Test listing versions when none exist"""
        response = await client.get(f"/api/projects/{sample_project.id}/versions")
        
        assert response.status_code == 200
        assert response.json() == []
    
    @pytest.mark.asyncio
    async def test_list_versions(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion
    ):
        """Test listing versions"""
        response = await client.get(f"/api/projects/{sample_project.id}/versions")
        
        assert response.status_code == 200
        data = response.json()
        
        assert len(data) == 1
        assert data[0]["version_number"] == 1
        assert data[0]["status"] == "draft"


class TestAudioSettings:
    """Tests for audio settings API"""
    
    @pytest.mark.asyncio
    async def test_get_audio_settings(
        self,
        client: AsyncClient,
        sample_project: Project
    ):
        """Test getting audio settings"""
        response = await client.get(f"/api/projects/{sample_project.id}/audio_settings")
        
        assert response.status_code == 200
        data = response.json()
        
        assert "background_music_enabled" in data
        assert "voice_gain_db" in data
        assert "ducking_strength" in data
        assert "voice_id" in data
    
    @pytest.mark.asyncio
    async def test_update_audio_settings(
        self,
        client: AsyncClient,
        sample_project: Project
    ):
        """Test updating audio settings"""
        response = await client.put(
            f"/api/projects/{sample_project.id}/audio_settings",
            json={
                "background_music_enabled": True,
                "voice_gain_db": 3.0,
                "ducking_strength": "strong",
                "voice_id": "test-voice-123",
            }
        )
        
        assert response.status_code == 200
        assert response.json()["status"] == "updated"
        
        # Verify changes
        response = await client.get(f"/api/projects/{sample_project.id}/audio_settings")
        data = response.json()
        
        assert data["background_music_enabled"] is True
        assert data["voice_gain_db"] == 3.0
        assert data["ducking_strength"] == "strong"
        assert data["voice_id"] == "test-voice-123"
    
    @pytest.mark.asyncio
    async def test_audio_settings_not_found(self, client: AsyncClient):
        """Test audio settings for non-existent project"""
        fake_id = uuid.uuid4()
        response = await client.get(f"/api/projects/{fake_id}/audio_settings")
        
        assert response.status_code == 404


class TestTranslationRules:
    """Tests for translation rules API"""
    
    @pytest.mark.asyncio
    async def test_get_translation_rules(
        self,
        client: AsyncClient,
        sample_project: Project
    ):
        """Test getting translation rules"""
        response = await client.get(f"/api/projects/{sample_project.id}/translation_rules")
        
        assert response.status_code == 200
        data = response.json()
        
        assert "do_not_translate" in data
        assert "preferred_translations" in data
        assert "style" in data
    
    @pytest.mark.asyncio
    async def test_update_translation_rules(
        self,
        client: AsyncClient,
        sample_project: Project
    ):
        """Test updating translation rules"""
        response = await client.put(
            f"/api/projects/{sample_project.id}/translation_rules",
            json={
                "do_not_translate": ["IFRS", "ESG", "KPI"],
                "preferred_translations": [
                    {"term": "revenue", "lang": "ru", "translation": "выручка"}
                ],
                "style": "friendly"
            }
        )
        
        assert response.status_code == 200
        
        # Verify changes
        response = await client.get(f"/api/projects/{sample_project.id}/translation_rules")
        data = response.json()
        
        assert data["do_not_translate"] == ["IFRS", "ESG", "KPI"]
        assert len(data["preferred_translations"]) == 1
        assert data["style"] == "friendly"


class TestPPTXUpload:
    """Tests for PPTX upload functionality"""
    
    @pytest.mark.asyncio
    async def test_upload_pptx_project_not_found(self, client: AsyncClient, tmp_path):
        """Test upload to non-existent project"""
        fake_id = uuid.uuid4()
        
        # Create a minimal pptx file
        from pptx import Presentation
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(pptx_path))
        
        with open(pptx_path, "rb") as f:
            response = await client.post(
                f"/api/projects/{fake_id}/upload_pptx",
                files={"file": ("test.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
            )
        
        assert response.status_code == 404
    
    @pytest.mark.asyncio
    async def test_upload_pptx_invalid_file_type(
        self,
        client: AsyncClient,
        sample_project: Project,
        tmp_path
    ):
        """Test upload with invalid file type"""
        # Create a fake txt file
        txt_path = tmp_path / "test.txt"
        txt_path.write_text("not a pptx")
        
        with open(txt_path, "rb") as f:
            response = await client.post(
                f"/api/projects/{sample_project.id}/upload_pptx",
                files={"file": ("test.txt", f, "text/plain")}
            )
        
        assert response.status_code == 400
        assert "invalid file type" in response.json()["detail"].lower()
    
    @pytest.mark.asyncio
    async def test_upload_pptx_success(
        self,
        client: AsyncClient,
        sample_project: Project,
        tmp_path
    ):
        """Test successful PPTX upload"""
        from pptx import Presentation
        
        # Create test PPTX
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(pptx_path))
        
        with open(pptx_path, "rb") as f:
            response = await client.post(
                f"/api/projects/{sample_project.id}/upload_pptx",
                files={"file": ("test.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                data={"comment": "Test upload"}
            )
        
        assert response.status_code == 200
        data = response.json()
        
        assert "version_id" in data
        assert data["version_number"] == 1
        assert data["status"] == "uploaded"


class TestConvertPPTX:
    """Tests for PPTX conversion endpoint"""
    
    @pytest.mark.asyncio
    async def test_convert_pptx_version_not_found(
        self,
        client: AsyncClient,
        sample_project: Project
    ):
        """Test conversion with non-existent version"""
        fake_version_id = uuid.uuid4()
        
        response = await client.post(
            f"/api/projects/{sample_project.id}/versions/{fake_version_id}/convert"
        )
        
        assert response.status_code == 404
    
    @pytest.mark.asyncio
    async def test_convert_pptx_no_file(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion
    ):
        """Test conversion when no PPTX file uploaded"""
        response = await client.post(
            f"/api/projects/{sample_project.id}/versions/{sample_version.id}/convert"
        )
        
        assert response.status_code == 400
        assert "no pptx" in response.json()["detail"].lower()
    
    @pytest.mark.asyncio
    async def test_convert_pptx_queues_task(
        self,
        client: AsyncClient,
        sample_project: Project,
        sample_version: ProjectVersion,
        db_session: AsyncSession
    ):
        """Test that conversion queues a Celery task"""
        # Set pptx_asset_path
        sample_version.pptx_asset_path = "/tmp/test.pptx"
        await db_session.commit()
        
        # Mock Celery task
        mock_task = MagicMock()
        mock_task.id = "test-task-123"
        
        with patch("app.workers.tasks.convert_pptx_task") as mock_convert:
            mock_convert.delay.return_value = mock_task
            
            response = await client.post(
                f"/api/projects/{sample_project.id}/versions/{sample_version.id}/convert"
            )
            
            assert response.status_code == 200
            data = response.json()
            
            assert data["task_id"] == "test-task-123"
            assert data["status"] == "queued"
            mock_convert.delay.assert_called_once()


class TestMusicUpload:
    """Tests for music file upload"""
    
    @pytest.mark.asyncio
    async def test_upload_music_invalid_format(
        self,
        client: AsyncClient,
        sample_project: Project,
        tmp_path
    ):
        """Test upload with invalid audio format"""
        wav_file = tmp_path / "music.wav"
        wav_file.write_bytes(b"fake wav data")
        
        with open(wav_file, "rb") as f:
            response = await client.post(
                f"/api/projects/{sample_project.id}/upload_music",
                files={"file": ("music.wav", f, "audio/wav")}
            )
        
        assert response.status_code == 400
        assert "mp3" in response.json()["detail"].lower()
    
    @pytest.mark.asyncio
    async def test_upload_music_success(
        self,
        client: AsyncClient,
        sample_project: Project,
        tmp_path
    ):
        """Test successful music upload"""
        mp3_file = tmp_path / "music.mp3"
        mp3_file.write_bytes(b"fake mp3 data")
        
        with open(mp3_file, "rb") as f:
            response = await client.post(
                f"/api/projects/{sample_project.id}/upload_music",
                files={"file": ("music.mp3", f, "audio/mpeg")}
            )
        
        assert response.status_code == 200
        data = response.json()
        
        assert "asset_id" in data
        assert data["status"] == "uploaded"


class TestVoices:
    """Tests for ElevenLabs voices endpoint"""

    @pytest.mark.asyncio
    async def test_get_voices_fallback_without_api_key(self, client: AsyncClient):
        """
        If ELEVENLABS_API_KEY is empty, endpoint should return a safe default voice
        without making any network calls.
        """
        # Reset module cache to avoid cross-test contamination
        import app.api.routes.projects as projects_routes
        projects_routes._voices_cache = {"voices": [], "timestamp": 0}

        with patch("app.api.routes.projects.settings.ELEVENLABS_API_KEY", ""):
            response = await client.get("/api/projects/voices")

        assert response.status_code == 200
        data = response.json()
        assert "voices" in data
        assert isinstance(data["voices"], list)
        assert len(data["voices"]) >= 1
        assert data["voices"][0]["voice_id"]

